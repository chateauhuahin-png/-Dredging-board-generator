"""
Board Builder - สร้างบอร์ดชี้แจงจาก PPTX + รูปภาพ
"""
import os, gc
from PIL import Image, ImageDraw, ImageFont, ImageFile
from pptx import Presentation

ImageFile.LOAD_TRUNCATED_IMAGES = True

BASE_DIR  = os.path.dirname(os.path.abspath(__file__))
FONT_BOLD = os.path.join(BASE_DIR, "fonts", "THSarabun Bold.ttf")
FONT_REG  = os.path.join(BASE_DIR, "fonts", "THSarabun.ttf")

# 150 DPI @ 120x80 cm
W, H   = 7087, 4724
NAVY   = (11, 20, 100)
GOLD   = (255, 215, 0)
WHITE  = (255, 255, 255)
MG, GAP = 120, 36
BW      = 6
UW = W - 2*MG
UH = H - 2*MG
HDR  = 400
PHH  = 840
LH   = 100
CL   = int(UW * 0.295)
CM   = int(UW * 0.325)
CR   = UW - CL - CM - 2*GAP
CONH = UH - HDR - GAP - PHH - GAP
CONY = MG + HDR + GAP
XL   = MG
XM   = XL + CL + GAP
XR   = XM + CM + GAP


def fnt(sz, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REG, sz)


def fit(path, w, h, bg=WHITE):
    """Contain mode"""
    canvas = Image.new("RGB", (w, h), bg)
    try:
        img = Image.open(path).convert("RGB")
        iw, ih = img.size
        sc = min(w / iw, h / ih)
        nw, nh = int(iw * sc), int(ih * sc)
        img = img.resize((nw, nh), Image.LANCZOS)
        canvas.paste(img, ((w - nw) // 2, (h - nh) // 2))
        img.close()
        del img
    except Exception as e:
        print(f"  fit error {path}: {e}")
    return canvas


def fit_cover(path, w, h):
    """Cover mode - scale to fill, crop edges"""
    canvas = Image.new("RGB", (w, h), WHITE)
    try:
        img = Image.open(path).convert("RGB")
        iw, ih = img.size
        sc = max(w / iw, h / ih)
        nw, nh = int(iw * sc), int(ih * sc)
        img = img.resize((nw, nh), Image.LANCZOS)
        ox = (nw - w) // 2
        oy = (nh - h) // 2
        canvas.paste(img.crop((ox, oy, ox + w, oy + h)), (0, 0))
        img.close()
        del img
    except Exception as e:
        print(f"  fit_cover error {path}: {e}")
    return canvas


def sec(draw, board, label, img_path, x, y, w, h, lsz=76):
    draw.rectangle([x, y, x+w, y+h], fill=WHITE)
    draw.rectangle([x, y, x+w, y+LH], fill=NAVY)
    f = fnt(lsz, bold=True)
    bb = draw.textbbox((0, 0), label, font=f)
    draw.text((x + (w-(bb[2]-bb[0]))//2, y + (LH-(bb[3]-bb[1]))//2),
              label, font=f, fill=GOLD)
    if img_path and os.path.exists(img_path):
        tile = fit(img_path, w, h - LH)
        board.paste(tile, (x, y + LH))
        tile.close()
        del tile
        gc.collect()
    draw.rectangle([x, y, x+w, y+h], outline=NAVY, width=BW)


def sec_multi(draw, board, label, img_paths, x, y, w, h, lsz=76):
    draw.rectangle([x, y, x+w, y+h], fill=WHITE)
    draw.rectangle([x, y, x+w, y+LH], fill=NAVY)
    f = fnt(lsz, bold=True)
    bb = draw.textbbox((0, 0), label, font=f)
    draw.text((x + (w-(bb[2]-bb[0]))//2, y + (LH-(bb[3]-bb[1]))//2),
              label, font=f, fill=GOLD)
    valid = [p for p in img_paths if p and os.path.exists(p)]
    if valid:
        n      = len(valid)
        slot_w = w // n
        img_h  = h - LH
        for i, p in enumerate(valid):
            tile = fit_cover(p, slot_w, img_h) if n == 1 else fit(p, slot_w, img_h)
            board.paste(tile, (x + i * slot_w, y + LH))
            tile.close()
            del tile
        gc.collect()
    draw.rectangle([x, y, x+w, y+h], outline=NAVY, width=BW)


def _parse_pptx_once(pptx_path, work_dir, med_dir):
    print("Opening PPTX (single pass)...")
    prs = Presentation(pptx_path)

    cfg = {"map": None, "surv": None, "des": None, "cross": None,
           "vol": None, "pr6": None, "pr4": None, "letter2": None}

    kw_detect = {
        "แผนที่":            "map",
        "ซ้ำซ้อน":           "letter2",
        "ตารางการสำรวจ":     "surv",
        "ตารางการออกแบบ":    "des",
        "รูปตัดตามขวาง":     "cross",
        "ตารางคำนวณปริมาตร": "vol",
        "แบบสรุป":           "pr6",
        "แบบประเมิน":        "pr4",
    }

    title1, title2, agency = "งานขุดลอกลำน้ำ", "", ""
    slide0 = prs.slides[0]
    for shape in slide0.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text = shape.text.strip()
        if not agency and ("นพค" in text or "นทพ" in text):
            for line in [l.strip() for l in text.split("\n") if l.strip()]:
                if "นพค" in line or "นทพ" in line:
                    agency = line; break
        if "งานขุด" in text:
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            title1 = f"{lines[0]} {lines[1]}" if len(lines) > 1 else lines[0]
            title2 = lines[2] if len(lines) > 2 else ""

    PHOTO_LABELS = {"ก่อน", "ระหว่าง", "หลัง"}
    photo_slide = None
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            continue
        text_all = " ".join(s.text for s in slide.shapes if hasattr(s, "text"))
        for kw, key in kw_detect.items():
            if kw in text_all and cfg[key] is None:
                cfg[key] = i
        if photo_slide is None and all(kw in text_all for kw in PHOTO_LABELS):
            photo_slide = slide

    print(f"Slide map: {cfg}")

    photo_before = photo_during = photo_after = None
    if photo_slide:
        label_x = {}
        for shape in photo_slide.shapes:
            if not hasattr(shape, "text"):
                continue
            for kw in PHOTO_LABELS:
                if kw in shape.text and kw not in label_x:
                    label_x[kw] = (shape.left or 0) + (shape.width or 0) // 2
        imgs = []
        for shape in photo_slide.shapes:
            if shape.shape_type == 13:
                cx = (shape.left or 0) + (shape.width or 0) // 2
                imgs.append((cx, shape.image.blob, shape.image.ext))
        imgs.sort(key=lambda x: x[0])
        if imgs:
            res = {}
            for kw, lx in label_x.items():
                cl = min(imgs, key=lambda img: abs(img[0] - lx))
                fname = os.path.join(work_dir, f"photo_{kw}.{cl[2]}")
                with open(fname, "wb") as f:
                    f.write(cl[1])
                res[kw] = fname
            photo_before = res.get("ก่อน")
            photo_during = res.get("ระหว่าง")
            photo_after  = res.get("หลัง")
        del imgs

    os.makedirs(med_dir, exist_ok=True)
    def _extract(shapes, prefix):
        pic_shapes = []
        for shape in shapes:
            if shape.shape_type == 13:
                pic_shapes.append(shape)
            elif shape.shape_type == 6:
                _extract(shape.shapes, f"{prefix}g")
        pic_shapes.sort(key=lambda s: (s.left or 0))
        for j, shape in enumerate(pic_shapes):
            img = shape.image
            fname = f"{prefix}_{j:03d}.{img.ext}"
            with open(os.path.join(med_dir, fname), "wb") as f:
                f.write(img.blob)

    key_slides = list(set(s for s in [cfg["map"], cfg["letter2"],
                                       cfg["surv"], cfg["des"], cfg["cross"],
                                       cfg["vol"], cfg["pr6"], cfg["pr4"]] if s))
    for si in key_slides:
        _extract(prs.slides[si-1].shapes, f"s{si:02d}")

    del prs
    gc.collect()
    print("PPTX closed, memory freed.")

    return cfg, title1, title2, agency, photo_before, photo_during, photo_after


def build_board(pptx_path, work_dir, output_path):
    logo_path = os.path.join(BASE_DIR, "fonts", "logo.png")
    os.makedirs(work_dir, exist_ok=True)
    med_dir = os.path.join(work_dir, "media")

    cfg, title1, title2, agency, photo_before, photo_during, photo_after = \
        _parse_pptx_once(pptx_path, work_dir, med_dir)

    def find_img(si):
        if si is None:
            return None
        candidates = []
        for f in sorted(os.listdir(med_dir)):
            if f.startswith(f"s{si:02d}_"):
                p = os.path.join(med_dir, f)
                try:
                    img = Image.open(p)
                    candidates.append((p, img.size[0]*img.size[1]))
                    img.close()
                except Exception:
                    pass
        if not candidates:
            return None
        candidates.sort(key=lambda x: x[1], reverse=True)
        return candidates[0][0]

    def find_all_imgs(si):
        if si is None:
            return []
        candidates = []
        for f in sorted(os.listdir(med_dir)):
            if f.startswith(f"s{si:02d}_"):
                p = os.path.join(med_dir, f)
                if os.path.exists(p):
                    candidates.append(p)
        return candidates

    map_jpg    = find_img(cfg["map"])
    pr6_img    = find_img(cfg["pr6"])
    pr4_img    = find_img(cfg["pr4"])
    surv_imgs  = find_all_imgs(cfg["surv"])
    des_imgs   = find_all_imgs(cfg["des"])
    cross_imgs = find_all_imgs(cfg["cross"])
    vol_imgs   = find_all_imgs(cfg["vol"])
    lett2_imgs = find_all_imgs(cfg["letter2"]) if cfg["letter2"] else []

    board = Image.new("RGB", (W, H), NAVY)
    draw  = ImageDraw.Draw(board)

    draw.rectangle([MG, MG, W-MG, MG+HDR], fill=NAVY)
    logo_sz = 680

    if logo_path and os.path.exists(logo_path):
        try:
            lg_img = Image.open(logo_path).convert("RGBA")
            lw, lh_img = lg_img.size
            sc = min(logo_sz / lw, (HDR - 20) / lh_img)
            nw, nh = int(lw * sc), int(lh_img * sc)
            lg_img = lg_img.resize((nw, nh), Image.LANCZOS)
            lx = XL + (logo_sz - nw) // 2
            ly_logo = MG + (HDR - nh) // 2
            board.paste(lg_img, (lx, ly_logo), mask=lg_img.split()[3])
        except Exception as e:
            print(f"  logo error: {e}")

    agency_text = agency if agency else ""
    fa = fnt(122, bold=True)
    bb_a = draw.textbbox((0,0), agency_text, font=fa)
    aw_txt = bb_a[2] - bb_a[0]
    ah_txt = bb_a[3] - bb_a[1]
    right_x = W - MG - logo_sz
    ax = right_x - aw_txt
    ay = MG + (HDR - ah_txt) // 2
    draw.text((ax, ay), agency_text, font=fa, fill=WHITE)

    t1sz = 148 if len(title1) < 60 else 132
    f1 = fnt(t1sz, True); f2 = fnt(114)
    cx = W // 2
    bb1 = draw.textbbox((0,0), title1, font=f1)
    bb2 = draw.textbbox((0,0), title2, font=f2)
    draw.text((cx-(bb1[2]-bb1[0])//2, MG+45),  title1, font=f1, fill=GOLD)
    draw.text((cx-(bb2[2]-bb2[0])//2, MG+205), title2, font=f2, fill=WHITE)

    boq_h = int(CONH * 0.54)
    sec(draw, board, "ประมาณการ (ปร.6)", pr6_img, XL, CONY, CL, boq_h)
    sec(draw, board, "ประมาณการ (ปร.4)", pr4_img, XL, CONY+boq_h+GAP, CL, CONH-boq_h-GAP)

    map_h = int(CONH * 0.50)
    sec(draw, board, "แผนที่และจุดดำเนินการ (มาตราส่วน 1:50,000)",
        map_jpg, XM, CONY, CM, map_h, lsz=68)
    sy = CONY + map_h + GAP
    sh = int(CONH * 0.265)
    sec_multi(draw, board, "ตารางการสำรวจ",  surv_imgs, XM, sy,        CM, sh)
    sec_multi(draw, board, "ตารางการออกแบบ", des_imgs,  XM, sy+sh+GAP, CM, CONH-map_h-GAP-sh-GAP)

    ch = int(CONH * 0.375)
    sec_multi(draw, board, "รูปตัดตามขวางขุดลอกลำน้ำ", cross_imgs, XR, CONY, CR, ch)
    ly = CONY + ch + GAP
    lh = int(CONH * 0.37)
    sec_multi(draw, board, "หนังสือตรวจสอบความซ้ำซ้อน", lett2_imgs, XR, ly, CR, lh, lsz=66)
    vy = ly + lh + GAP
    sec_multi(draw, board, "ตารางคำนวณปริมาตรดินตะกอน", vol_imgs, XR, vy, CR, CONH-ch-GAP-lh-GAP)

    phy = CONY + CONH + GAP
    PW  = (UW - 2*GAP) // 3
    for idx, (lbl, ph) in enumerate(zip(
        ["ภาพก่อนปฏิบัติงาน", "ภาพระหว่างปฏิบัติงาน", "ภาพหลังปฏิบัติงาน"],
        [photo_before, photo_during, photo_after]
    )):
        px = XL + idx * (PW + GAP)
        draw.rectangle([px, phy, px+PW, phy+PHH], fill=WHITE)
        draw.rectangle([px, phy, px+PW, phy+LH], fill=NAVY)
        fp = fnt(84, True)
        bbl = draw.textbbox((0,0), lbl, font=fp)
        draw.text((px+(PW-(bbl[2]-bbl[0]))//2, phy+(LH-(bbl[3]-bbl[1]))//2),
                  lbl, font=fp, fill=GOLD)
        if ph and os.path.exists(ph):
            board.paste(fit_cover(ph, PW, PHH-LH), (px, phy+LH))
        draw.rectangle([px, phy, px+PW, phy+PHH], outline=NAVY, width=BW)

    board.save(output_path, "JPEG", quality=95, dpi=(150, 150))
    print(f"Saved: {output_path}")
    return output_path
